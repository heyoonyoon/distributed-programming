#!/usr/bin/env python3
# 절제된 에디토리얼: 흰 캔버스 + 거의 검정 + 인디고 액센트 한 색(+톤) + 다크 네이비.
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE=os.path.dirname(os.path.abspath(__file__)); DIA=os.path.join(HERE,"diagrams"); CL=os.path.join(DIA,"clusters")

def C(h): return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
BLACK=C("1a1a1a"); WHITE=C("ffffff"); HAIR=C("e6e6e6"); SOFT=C("f4f4f6"); GRAY=C("6b7280")
NAVY=C("1f1d3d")
ACCENT=C("5b5bd6"); ACCENT_SOFT=C("ecebfb"); ACCENT_MED=C("c9c4f3"); ACCENT_LT=C("cdbff7")
WARM=C("f7ede6")
SANS="Apple SD Gothic Neo"; MONO="Menlo"
A_CUST=ACCENT_SOFT; A_STAFF=ACCENT; A_SYS=NAVY      # 인디고 명도 3단계
def fg(fill): return WHITE if fill in (ACCENT,NAVY) else BLACK

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW=13.333

def slide(bg=WHITE):
    s=prs.slides.add_slide(BLANK); f=s.background.fill; f.solid(); f.fore_color.rgb=bg; return s

def _runs(tf,runs,align,anchor,wrap,spc):
    tf.word_wrap=wrap; tf.vertical_anchor=anchor; first=True
    for line in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False; p.alignment=align
        for (t,sz,b,c,nm) in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name=nm or SANS
            if spc is not None: r.font._rPr.set('spc',str(int(spc*100)))

def tb(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,wrap=True,spc=None):
    x=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); _runs(x.text_frame,runs,align,anchor,wrap,spc); return x

def panel(s,l,t,w,h,fill,rad=0.12,line=None):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h)); sp.shadow.inherit=False
    try: sp.adjustments[0]=rad
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1.0)
    return sp

def box(s,l,t,w,h,fill,runs,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,rad=0.16,line=None,wrap=True,shape=MSO_SHAPE.ROUNDED_RECTANGLE,spc=None):
    sp=s.shapes.add_shape(shape,Inches(l),Inches(t),Inches(w),Inches(h)); sp.shadow.inherit=False
    if shape==MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0]=rad
        except Exception: pass
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1.2)
    tf=sp.text_frame; tf.margin_top=Pt(6); tf.margin_bottom=Pt(6); tf.margin_left=Pt(14); tf.margin_right=Pt(14)
    _runs(tf,runs,align,anchor,wrap,spc); return sp

def head(s,eyebrow,title,ink=BLACK):
    # eyebrow(영어 소제목)는 쓰지 않는다. 제목만 크게.
    tb(s,0.82,0.62,11.9,1.5,[[(title,52,True,ink,SANS)]],spc=-1.0)
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.9),Inches(1.78),Inches(0.95),Inches(0.09)); bar.shadow.inherit=False
    bar.fill.solid(); bar.fill.fore_color.rgb=ACCENT; bar.line.fill.background()

def arrow(s,l,t,sz=30,col=ACCENT):
    tb(s,l,t,0.7,0.7,[[("→",sz,True,col,SANS)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

def img_fit(s,path,bl,bt,bw,bh):
    iw,ih=Image.open(path).size; r=iw/ih; w,h=bw,bw/r
    if h>bh: h,w=bh,bh*r
    s.shapes.add_picture(path,Inches(bl+(bw-w)/2),Inches(bt+(bh-h)/2),Inches(w),Inches(h))

# 1. 타이틀
s=slide(NAVY)
tb(s,0.95,2.35,11.6,2.2,[[("설계대로 구현하고,",56,True,WHITE,SANS)],[("AI를 통제하다",56,True,WHITE,SANS)]],spc=-1.2)
tb(s,1.0,4.85,11.3,0.9,[[("의료·자동차 보험 시스템",30,False,ACCENT_LT,SANS)]],spc=-0.3)

# 2. 세 주체
s=slide(); head(s,"ACTORS","누가 사용하나")
data=[("고객",A_CUST,"가입·청구·납부"),("직원",A_STAFF,"가입·지급 심사"),("시스템",A_SYS,"계약·고지서·지급 자동화")]
xs=[1.55,5.42,9.29]
for (lab,col,role),x in zip(data,xs):
    box(s,x,2.55,2.5,2.5,col,[[(lab,34,True,fg(col),SANS)]],shape=MSO_SHAPE.OVAL,wrap=False)
    tb(s,x-1.0,5.3,4.5,0.8,[[(role,27,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,spc=-0.3)

# 3. 두 가지 보험
s=slide(); head(s,"PRODUCTS","두 가지 보험")
for nm,ess,col,x in [("의료보험","아플 때 — 병원비",ACCENT_SOFT,1.45),("자동차보험","사고 났을 때 — 수리·치료비",WARM,7.0)]:
    panel(s,x,2.6,4.9,3.1,col,rad=0.12)
    tb(s,x,3.4,4.9,1.0,[[(nm,42,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,spc=-0.6)
    tb(s,x,4.65,4.9,0.8,[[(ess,26,False,GRAY,SANS)]],align=PP_ALIGN.CENTER,spc=-0.3)

# 4. 의료보험 — 분기
s=slide(); head(s,"MEDICAL","병원비를 보험금으로 돌려준다")
box(s,0.9,3.95,3.1,1.4,ACCENT,[[("병원비 청구",28,True,WHITE,SANS)]],wrap=False)
arrow(s,4.25,3.15,32); arrow(s,4.25,5.15,32)
box(s,5.1,2.75,7.4,1.5,ACCENT_SOFT,[[("100만 원 미만 → 바로 입금",27,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,wrap=False)
box(s,5.1,4.7,7.4,1.5,WARM,[[("100만 원 이상 → 직원 확인 후 입금",27,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,wrap=False)

# 5. 자동차보험 — 흐름
s=slide(); head(s,"AUTO","사고 피해를 보상한다")
steps=["사고 접수","직원이 맡음","보상액 결정","계좌 입금"]
iw,gap=2.6,0.55; total=len(steps)*iw+(len(steps)-1)*gap; x=(SW-total)/2
for i,st in enumerate(steps):
    col=ACCENT_SOFT if i==len(steps)-1 else SOFT
    box(s,x,3.7,iw,1.5,col,[[(st,26,True,BLACK,SANS)]],line=HAIR,wrap=False)
    if i<len(steps)-1: arrow(s,x+iw+0.0,3.95,30)
    x+=iw+gap

# 6-8. 여정
def journey(eye,t,steps,vtext):
    s=slide(); head(s,eye,t); n=len(steps); iw,gap=2.45,0.7
    total=n*iw+(n-1)*gap; x=(SW-total)/2
    cmap={"cust":A_CUST,"staff":A_STAFF,"sys":A_SYS}
    for i,(who,lab,act) in enumerate(steps):
        col=cmap[who]
        box(s,x+iw/2-0.85,2.5,1.7,1.7,col,[[(lab,24,True,fg(col),SANS)]],shape=MSO_SHAPE.OVAL,wrap=False)
        box(s,x,4.5,iw,1.3,col,[[(act,26,True,fg(col),SANS)]],wrap=False)
        if i<n-1: arrow(s,x+iw+0.05,4.75,30)
        x+=iw+gap
    tb(s,0.85,6.15,11.6,1.0,[[(vtext,27,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,spc=-0.4)

journey("JOURNEY 1","가입",
    [("cust","고객","상품 조회"),("cust","고객","가입 신청"),("staff","직원","가입 심사"),("sys","시스템","계약 생성")],
    "승인되면 계약이 바로 생깁니다.")
journey("JOURNEY 2","납부와 미납",
    [("cust","고객","보험료 납부"),("sys","시스템","미납 감지"),("sys","시스템","고지서 발송")],
    "납부가 밀리면 미납으로 잡히고, 매일 아침 고지서가 나갑니다 (30일 초과 시 해지 예고).")
journey("JOURNEY 3","보상과 지급",
    [("cust","고객","청구·접수"),("staff","직원","보상 심사"),("sys","시스템","보험금 지급")],
    "병원비는 100만 원 미만이면 즉시, 이상이거나 자동차 사고는 직원 심사 후 지급.")

# 9. 도메인 여섯 개
s=slide(); head(s,"DOMAINS","도메인 여섯 개")
doms=["사용자","상품","계약","청구","심사","사고이력"]
cw,gap=1.95,0.26; total=len(doms)*cw+(len(doms)-1)*gap; x=(SW-total)/2
for d in doms:
    box(s,x,3.25,cw,1.45,SOFT,[[(d,27,True,BLACK,SANS)]],line=HAIR,wrap=False); x+=cw+gap

# 10-16. 비교
DC=os.path.join(DIA,"design-crops")
comps=[("사용자","user","글자로 된 번호를 숫자로 바꾸고, 날짜 표기만 손봤다."),
       ("상품","product","연결 구조(계약·보장 항목·하위 상품)는 그대로 4개. 번호를 숫자로 바꾸고, 보험료 계산은 아직 뼈대만 있다."),
       ("계약","contract","자동이체 정보를 더했다."),
       ("청구","claim","첨부파일을 더하고, 상태를 4개 → 6개로 늘려 송금 결과까지 추적한다."),
       ("__CODE__",None,None),
       ("심사","review","지급 심사가 의료뿐 아니라 자동차 사고까지 맡도록 넓혔다."),
       ("사고이력","accident","외부 연동이 아직 가짜라, 사고 한 건 클래스는 빼고 단순화했다.")]
for name,key,diff in comps:
    if name=="__CODE__":
        s=slide(); head(s,"","청구 상태값")
        code=("public enum ClaimStatus {\n    PENDING, IN_REVIEW, APPROVED, REJECTED,   // 설계 4개\n"
              "    COMPLETED, FAILED                         // 구현에서 추가: 송금 성공·실패\n}")
        panel(s,0.95,2.6,11.4,2.1,SOFT,rad=0.04)
        tb(s,1.3,2.8,10.8,1.8,[[(l,19,False,BLACK,MONO)] for l in code.split("\n")],anchor=MSO_ANCHOR.MIDDLE)
        tb(s,0.85,5.05,11.6,1.0,[[("돈을 실제로 보냈는지까지 따라가려고 상태 두 개를 더했다.",28,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,spc=-0.4); continue
    s=slide(); head(s,"",name)
    tb(s,0.5,1.95,6.1,0.55,[[("내가 그린 설계",24,True,ACCENT,SANS)]],align=PP_ALIGN.CENTER,spc=-0.3)
    tb(s,6.75,1.95,6.1,0.55,[[("실제 코드",24,True,C("b5793a"),SANS)]],align=PP_ALIGN.CENTER,spc=-0.3)
    img_fit(s,os.path.join(CL,f"{key}-design.png"),0.5,2.5,6.1,3.35)
    img_fit(s,os.path.join(CL,f"{key}-code.png"),6.75,2.5,6.1,3.35)
    box(s,0.5,6.05,12.35,1.15,WARM,[[("다른 점   ",19,True,C("b5793a"),SANS),(diff,19,True,BLACK,SANS)]],align=PP_ALIGN.LEFT,rad=0.12)

# 17. 유스케이스
s=slide(); head(s,"","전체 중 내가 시연하는 부분")
img_fit(s,os.path.join(DIA,"uc-map.png"),4.6,2.0,4.1,4.1)
box(s,0.85,2.3,3.4,0.8,ACCENT,[[("파란 박스 = 직접 시연",18,True,WHITE,SANS)]],wrap=False)
box(s,0.85,3.3,3.4,0.8,C("f1f1f3"),[[("회색 = 시연 외(코드엔 있음)",17,True,GRAY,SANS)]],line=C("c9c9cf"),wrap=False)
tb(s,0.85,6.35,11.6,0.8,[[("전체 흐름 중 파란 부분만 라이브로 시연한다.",26,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,spc=-0.4)

# ===== 프로젝트 구조 =====
s=slide(); head(s,"","프로젝트 구조")
def varrow(s,l,t):
    tb(s,l,t,0.9,0.32,[[("↓",18,True,GRAY,SANS)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
# 백엔드 — 계층형
tb(s,0.85,2.0,5.9,0.5,[[("백엔드 · Spring — 계층형",22,True,ACCENT,SANS)]],spc=-0.3)
blayers=["controller","service","domain","repository","DB"]
y=2.6
for i,nm in enumerate(blayers):
    fillc=ACCENT if nm=="DB" else (ACCENT_SOFT if nm=="domain" else SOFT)
    box(s,1.85,y,3.9,0.56,fillc,[[(nm,18,True,fg(fillc) if nm=="DB" else BLACK,MONO)]],line=HAIR,wrap=False)
    if i<len(blayers)-1: varrow(s,3.35,y+0.5)
    y+=0.7
tb(s,0.85,y+0.0,5.9,0.45,[[("dto = 요청·응답 데이터 (controller ↔ 화면)",14,False,GRAY,MONO)]],spc=-0.2,wrap=False)
# 프론트 — 기능형
tb(s,7.05,2.0,5.5,0.5,[[("프론트 · React — 기능형",22,True,C("b5793a"),SANS)]],spc=-0.3)
box(s,7.05,2.65,5.5,2.4,SOFT,[[("features/",19,True,BLACK,MONO)],
    [("기능 한 폴더 = 화면 + hook + api",15,False,GRAY,SANS)],[("",8,False,GRAY,SANS)],
    [("auth · contracts · claims",16,False,BLACK,SANS)],
    [("benefit-review · underwriting-review",16,False,BLACK,SANS)],
    [("assignments · profile · customer-home",16,False,BLACK,SANS)]],align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE,line=HAIR)
tb(s,0.85,6.65,11.6,0.7,[[("백엔드는 계층으로, 프론트는 기능 단위로 나눴다.",26,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,spc=-0.4)

# ===== 코드 워크스루: 상품 상세 조회 =====
SHOTS=os.path.join(DIA,"shots")
LILAC=C("cdbff7"); CMNT=C("7c83b0")
def kw(l):
    ls=l.strip()
    if ls.startswith("//") or ls.startswith("#"): return CMNT
    if any(k in l for k in ("@","class ","public ","return","async ","function","const ","export","import","void","static ","interface","SELECT")): return LILAC
    return WHITE
def codepage(step,title,where,lines,caption):
    s=slide(); head(s,"",title)
    tb(s,0.9,1.95,11.6,0.5,[[(step+"   ",18,True,ACCENT,SANS),(where,16,False,GRAY,MONO)]],spc=-0.2)
    panel(s,0.9,2.5,11.55,3.45,NAVY,rad=0.04)
    tb(s,1.3,2.7,10.9,3.05,[[(l,15,False,kw(l),MONO)] for l in lines],anchor=MSO_ANCHOR.MIDDLE)
    tb(s,0.9,6.15,11.6,0.8,[[(caption,24,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,spc=-0.4)

# 시나리오 + BEFORE
s=slide(); head(s,"","따라가 볼 한 번의 클릭")
tb(s,0.9,1.95,11.6,0.6,[[("고객이 상품 목록에서 ",20,True,BLACK,SANS),("'실손기본'",20,True,ACCENT,SANS),("을 클릭한다.",20,True,BLACK,SANS)]],spc=-0.3)
img_fit(s,os.path.join(SHOTS,"products-before.png"),1.6,2.5,10.1,3.5)
tb(s,0.9,6.2,11.6,0.8,[[("이 클릭 하나가 코드로 어떻게 흘러 화면까지 오는지 따라가 본다.",24,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,spc=-0.4)

codepage("①","프론트 — 버튼 클릭","features/contracts/ProductsView.tsx",
 ["// 상품 목록의 각 카드","<button onClick={() => selectProduct(product.id)}>",
  "  {product.productName}","</button>","",
  "// 클릭하면 실행되는 함수","async function selectProduct(id) {",
  "  setSelectedProduct(await contractsApi.getProduct(id))","}"],
 "카드를 누르면 그 상품 id로 selectProduct가 실행된다.")

codepage("②","프론트 — API 호출","contracts/api.ts · lib/api/httpClient.ts",
 ["// api.ts","getProduct(id) {","  return request(`/products/${id}`)","}","",
  "// httpClient.ts","fetch(`${baseUrl}${path}`)  // baseUrl = http://localhost:8080",
  "// → GET http://localhost:8080/products/2"],
 "프론트가 백엔드의 /products/2 로 GET 요청을 보낸다.")

codepage("③","백엔드 — 컨트롤러","web/controller/ProductController.java",
 ["@RestController","class ProductController {","",
  '  @GetMapping("/products/{id}")',
  "  ProductDetailResponse detail(@PathVariable Long id) {",
  "    return productService.detail(id);","  }","}"],
 "컨트롤러가 요청을 받아 서비스로 넘긴다.")

codepage("④","백엔드 — 서비스 · 저장소","service/ProductService.java · repository/ProductRepository.java",
 ["@Transactional(readOnly = true)","ProductDetailResponse detail(Long id) {",
  "  return productRepository.findById(id)        // DB 조회",
  "      .map(ProductDetailResponse::from)",
  '      .orElseThrow(() -> new IllegalArgumentException("상품을 찾을 수 없습니다."));',"}","",
  "// ProductRepository extends JpaRepository<InsuranceProduct, Long>",
  "// findById → SELECT * FROM insurance_product WHERE id = 2"],
 "서비스가 저장소로 DB에서 상품을 꺼낸다.")

codepage("⑤","백엔드 — 엔티티 → 응답 DTO","web/dto/ProductDetailResponse.java",
 ["// DB 엔티티를 화면용 데이터로 변환",
  "static ProductDetailResponse from(InsuranceProduct product) {",
  "  return new ProductDetailResponse(",
  "    product.getId(), product.getProductName(),",
  "    product.getBasePremium(),",
  "    product.getCoverageItems().stream()",
  "        .map(CoverageItemResponse::from).toList());","}"],
 "DB 엔티티를 그대로 노출하지 않고 응답 DTO로 변환한다.")

# 응답 → 화면 + AFTER
s=slide(); head(s,"","⑥  응답이 화면에 표시된다")
panel(s,0.85,1.9,6.05,2.0,NAVY,rad=0.04)
tb(s,1.15,2.05,5.5,1.8,[[("// 백엔드 응답 (JSON)",13,False,CMNT,MONO)],
  [('{ "productName": "실손기본",',13,False,WHITE,MONO)],
  [('  "monthlyPremium": 12000,',13,False,WHITE,MONO)],
  [('  "coverageItems": [',13,False,WHITE,MONO)],
  [('    {"itemName":"통원치료비" …},',13,False,LILAC,MONO)],
  [('    {"itemName":"입원치료비" …} ]}',13,False,LILAC,MONO)]],anchor=MSO_ANCHOR.MIDDLE)
img_fit(s,os.path.join(SHOTS,"products-after.png"),7.1,1.9,5.5,4.0)
tb(s,0.85,5.05,6.05,1.0,[[("받은 데이터를 그대로 화면에",17,True,BLACK,SANS)],[("상세·보장항목으로 렌더한다.",17,True,BLACK,SANS)]],anchor=MSO_ANCHOR.MIDDLE)
tb(s,0.9,6.35,11.6,0.7,[[("클릭 → 컨트롤러 → 서비스 → DB → DTO → 화면, 한 줄기로 흐른다.",24,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,spc=-0.4)

# 18. AI 사용
s=slide(); head(s,"","AI를 어떻게 썼나")
panel(s,0.85,2.2,5.7,1.7,ACCENT_SOFT,rad=0.12)
tb(s,1.2,2.4,5.1,1.4,[[("사람",26,True,ACCENT,SANS)],[("설계 · 결정 · 검토",22,True,BLACK,SANS)]],anchor=MSO_ANCHOR.MIDDLE)
panel(s,6.75,2.2,5.7,1.7,SOFT,rad=0.12,line=HAIR)
tb(s,7.1,2.4,5.1,1.4,[[("AI",26,True,BLACK,SANS)],[("구현 · 테스트 · 정리",22,True,BLACK,SANS)]],anchor=MSO_ANCHOR.MIDDLE)
steps=["설계","따져묻기","계획","테스트 먼저 구현","검토"]; cw,gap=2.15,0.42; total=len(steps)*cw+(len(steps)-1)*gap; x=(SW-total)/2
for i,st in enumerate(steps):
    col=ACCENT_SOFT if st=="테스트 먼저 구현" else SOFT
    box(s,x,4.35,cw,1.2,col,[[(st,21,True,BLACK,SANS)]],line=HAIR,wrap=True)
    if i<len(steps)-1: arrow(s,x+cw+0.0,4.6,26)
    x+=cw+gap
tb(s,0.85,5.95,11.6,0.9,[[("설계와 결정은 사람이, 구현은 AI가 했다.",28,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,spc=-0.4)

# 19. AI 문제
s=slide(); head(s,"","AI가 삑사리 난 곳, 어떻게 잡았나")
cards=[("성급하게 쪼갬","쓸모없는 계층 분리 → 되돌림"),("시킨 것과 다름","자동 배정 → 수동 배정으로 정정"),
       ("정해둔 틀에서 이탈","점검으로 찾아내 관리"),("용어를 섞어 씀","용어집으로 미리 차단")]
pos=[(0.85,2.25),(6.75,2.25),(0.85,4.15),(6.75,4.15)]
for (h,fix),(x,y) in zip(cards,pos):
    panel(s,x,y,5.7,1.7,SOFT,rad=0.1,line=HAIR)
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y+0.25),Inches(0.12),Inches(1.2)); bar.shadow.inherit=False
    bar.fill.solid(); bar.fill.fore_color.rgb=ACCENT; bar.line.fill.background()
    tb(s,x+0.5,y+0.22,5.0,1.35,[[(h,24,True,BLACK,SANS)],[(fix,19,True,GRAY,SANS)]],anchor=MSO_ANCHOR.MIDDLE,spc=-0.3)
tb(s,0.85,6.15,11.6,0.9,[[("AI는 빠르지만, 방향은 사람이 잡았다.",28,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,spc=-0.4)

# 20. 결론
s=slide(NAVY)
tb(s,0.95,2.75,11.6,2.2,[[("설계한 그대로 구현했고,",48,True,WHITE,SANS)],[("그 과정에서 AI를 통제했다.",48,True,WHITE,SANS)]],spc=-1.0)

out=os.path.join(HERE,"보험발표-figma.pptx"); prs.save(out); print("saved:",out,"| slides:",len(prs.slides._sldIdLst))
