#!/usr/bin/env python3
# AI 도구 정리 / cmux / 왜 Codex 3페이지 삽입 (in-place).
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE=os.path.dirname(os.path.abspath(__file__)); PPTX=os.path.join(HERE,"보험발표-figma.pptx")
def C(h): return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
BLACK=C("1a1a1a"); WHITE=C("ffffff"); HAIR=C("e6e6e6"); SOFT=C("f4f4f6"); GRAY=C("6b7280")
NAVY=C("1f1d3d"); ACCENT=C("5b5bd6"); ACCENT_SOFT=C("ecebfb"); WARM=C("f7ede6"); BROWN=C("b5793a")
SANS="Apple SD Gothic Neo"; MONO="Menlo"; SW=13.333
def fg(fill): return WHITE if fill in (ACCENT,NAVY) else BLACK
prs=Presentation(PPTX); BLANK=prs.slide_layouts[6]
def newslide(bg=WHITE):
    s=prs.slides.add_slide(BLANK); f=s.background.fill; f.solid(); f.fore_color.rgb=bg; return s
def _runs(tf,runs,align,anchor,wrap,spc):
    tf.word_wrap=wrap; tf.vertical_anchor=anchor; first=True
    for line in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False; p.alignment=align
        for (t,sz,b,c,nm) in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name=nm or SANS
            if spc is not None: r.font._rPr.set('spc',str(int(spc*100)))
def tb(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,wrap=True,spc=None):
    x=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); _runs(x.text_frame,runs,align,anchor,wrap,spc); return x
def box(s,l,t,w,h,fill,runs,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,rad=0.12,line=None,wrap=True,spc=None):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h)); sp.shadow.inherit=False
    try: sp.adjustments[0]=rad
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1.3)
    tf=sp.text_frame; tf.margin_top=Pt(6); tf.margin_bottom=Pt(6); tf.margin_left=Pt(16); tf.margin_right=Pt(14)
    _runs(tf,runs,align,anchor,wrap,spc); return sp
def head(s,title):
    tb(s,0.82,0.62,11.9,1.5,[[(title,46,True,BLACK,SANS)]],spc=-1.0)
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.9),Inches(1.7),Inches(0.95),Inches(0.09)); bar.shadow.inherit=False
    bar.fill.solid(); bar.fill.fore_color.rgb=ACCENT; bar.line.fill.background()
def arrow(s,l,t,sz=34,col=ACCENT):
    tb(s,l,t,0.8,0.7,[[("→",sz,True,col,SANS)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

made=[]
# 1) 내가 쓴 AI 도구
s=newslide(); head(s,"내가 쓴 AI 도구"); made.append(s)
rows=[("Claude Code","설계·조율은 Opus, 구현은 Sonnet 서브에이전트"),
      ("Superpowers 스킬","brainstorming · grill-with-docs · writing/executing-plans"),
      ("Codex","다른 모델로 교차 리뷰"),
      ("cmux","프론트·백엔드 Claude를 메시지로 연결")]
y=2.3
for nm,desc in rows:
    box(s,0.9,y,11.55,0.95,SOFT,[[(nm+"     ",19,True,ACCENT,MONO),(desc,18,False,BLACK,SANS)]],align=PP_ALIGN.LEFT,line=HAIR,wrap=False)
    y+=1.12

# 2) cmux 하이라이트
s=newslide(); head(s,"cmux — 두 Claude를 붙이다"); made.append(s)
box(s,1.4,2.95,4.4,1.5,ACCENT_SOFT,[[("프론트 Claude",20,True,ACCENT,SANS)],[("React 화면",16,False,GRAY,SANS)]],anchor=MSO_ANCHOR.MIDDLE)
box(s,7.55,2.95,4.4,1.5,ACCENT_SOFT,[[("백엔드 Claude",20,True,ACCENT,SANS)],[("Spring API",16,False,GRAY,SANS)]],anchor=MSO_ANCHOR.MIDDLE)
tb(s,5.75,2.85,1.85,0.8,[[("⇄",30,True,ACCENT,SANS)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
tb(s,1.4,3.55,10.5,0.5,[[("cmux",14,True,GRAY,MONO)]],align=PP_ALIGN.CENTER,spc=0.4)
tb(s,0.9,5.0,11.6,1.2,[[("두 Claude 세션을 cmux로 메시지 교환시켜",24,True,BLACK,SANS)],[("한쪽이 API를 바꾸면 다른 쪽이 받아서 맞췄다.",24,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spc=-0.4)

# 3) 왜 Codex 교차 리뷰
s=newslide(); head(s,"왜 Codex로 교차 리뷰했나"); made.append(s)
box(s,0.85,3.0,5.25,2.1,SOFT,[[("같은 모델끼리 리뷰하면",18,True,GRAY,SANS)],[("",6,False,GRAY,SANS)],[("같은 맹점을 똑같이 놓친다",21,True,BLACK,SANS)]],align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE,line=HAIR)
arrow(s,6.27,3.7)
box(s,7.25,3.0,5.25,2.1,ACCENT_SOFT,[[("그래서",18,True,ACCENT,SANS)],[("",6,False,GRAY,SANS)],[("다른 모델(Codex)로 교차 검증",21,True,BLACK,SANS)]],align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE)
tb(s,0.9,5.6,11.6,0.8,[[("마지막에 Codex로 한 번 더 봐서 Claude가 놓친 것을 잡았다.",22,False,GRAY,SANS)]],align=PP_ALIGN.CENTER,spc=-0.3)

# ===== 위치 이동 =====
RID=qn('r:id'); sldIdLst=prs.slides._sldIdLst
id2el={int(el.get('id')):el for el in list(sldIdLst)}
def title_of(sl):
    for sh in sl.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip(): return sh.text_frame.text.strip().split("\n")[0]
    return ""
tools_el,cmux_el,codex_el=[id2el[s.slide_id] for s in made]
for el in (tools_el,cmux_el,codex_el): sldIdLst.remove(el)
def idx_after(prefix):
    for i,el in enumerate(list(sldIdLst)):
        sid=int(el.get('id'))
        sl=next(s for s in prs.slides if s.slide_id==sid)
        if title_of(sl).startswith(prefix): return i+1
    return len(list(sldIdLst))
# 도구·cmux: 분업 뒤
p=idx_after("사람과 AI"); sldIdLst.insert(p,tools_el); sldIdLst.insert(p+1,cmux_el)
# codex: 다중 리뷰 뒤
p=idx_after("다중 리뷰"); sldIdLst.insert(p,codex_el)
prs.save(PPTX)
print("도구/cmux/Codex 3페이지 삽입 완료 | total:", len(prs.slides._sldIdLst))
