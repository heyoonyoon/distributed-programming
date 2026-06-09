import sys,os,copy
from pptx import Presentation
src=sys.argv[1]; outdir=sys.argv[2]; os.makedirs(outdir,exist_ok=True)
base=Presentation(src); n=len(base.slides._sldIdLst)
for i in range(n):
    p=Presentation(src); lst=p.slides._sldIdLst; ids=list(lst)
    for j,el in enumerate(ids):
        if j!=i: lst.remove(el)
    tmp=os.path.join(outdir,f"s{i:02d}.pptx"); p.save(tmp)
print(n)
