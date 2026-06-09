#!/usr/bin/env python3
# 편집 가능한 네이티브 PPTX. 폰트·도형 크게, 색 또렷하게.
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
DIA = os.path.join(HERE, "diagrams"); CL = os.path.join(DIA, "clusters")

NAVY=RGBColor(0x0B,0x3D,0x5C); NAVY2=RGBColor(0x0F,0x2A,0x43); AMBER=RGBColor(0xF2,0xA9,0x00)
INK=RGBColor(0x1C,0x2B,0x36); MUTED=RGBColor(0x5B,0x71,0x85); LINE=RGBColor(0xC3,0xD2,0xDF)
WHITE=RGBColor(0xFF,0xFF,0xFF); AMBER_DK=RGBColor(0xB5,0x79,0x00)
# 액터/주체 색 (빨강·파랑·노랑 — 완전히 다른 색)
CUST=RGBColor(0x1E,0x64,0xC8); STAFF=RGBColor(0xD8,0x39,0x2B); SYS=RGBColor(0xF4,0xB4,0x00)
CUST_BG=RGBColor(0xDC,0xE7,0xF8); STAFF_BG=RGBColor(0xF8,0xDE,0xDB); SYS_BG=RGBColor(0xFB,0xEF,0xC4)
def actc(fill):  # 노랑 위엔 진한 글자, 그 외엔 흰 글자
    return NAVY if fill==SYS else WHITE
# 같음/다름
OKBG=RGBColor(0xDC,0xF1,0xEC); OK=RGBColor(0x15,0x80,0x73)
AMBERBG=RGBColor(0xFD,0xEC,0xC4)
FONT="Apple SD Gothic Neo"; MONO="Menlo"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW=13.333

def slide(bg=None):
    s=prs.slides.add_slide(BLANK)
    if bg is not None:
        f=s.background.fill; f.solid(); f.fore_color.rgb=bg
    return s

def _set(tf,runs,align,anchor,wrap=True):
    tf.word_wrap=wrap; tf.vertical_anchor=anchor; first=True
    for line in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False; p.alignment=align
        for (t,sz,b,c,nm) in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b
            r.font.color.rgb=c; r.font.name=nm or FONT

def tb(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    x=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); _set(x.text_frame,runs,align,anchor); return x

def box(s,l,t,w,h,fill,line,runs,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,lw=1.5,shape=MSO_SHAPE.ROUNDED_RECTANGLE,wrap=True):
    sp=s.shapes.add_shape(shape,Inches(l),Inches(t),Inches(w),Inches(h)); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    tf=sp.text_frame; tf.margin_top=Pt(5); tf.margin_bottom=Pt(5); tf.margin_left=Pt(10); tf.margin_right=Pt(10)
    _set(tf,runs,align,anchor,wrap); return sp

def bar(s,l,t,w,col,h=0.14):
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h)); r.shadow.inherit=False
    r.fill.solid(); r.fill.fore_color.rgb=col; r.line.fill.background()

def title(s,text):
    tb(s,0.6,0.4,12.1,1.0,[[(text,36,True,NAVY,FONT)]],anchor=MSO_ANCHOR.MIDDLE)
    bar(s,0.62,1.42,2.9,AMBER,0.07)

def arrow(s,l,t,sz=36):
    tb(s,l,t,0.7,0.7,[[("→",sz,True,AMBER,FONT)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

def verdict(s,text,t=6.1):
    tb(s,0.6,t,12.13,1.0,[[(text,24,True,NAVY,FONT)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

def img_fit(s,path,bl,bt,bw,bh):
    iw,ih=Image.open(path).size; r=iw/ih; w,h=bw,bw/r
    if h>bh: h,w=bh,bh*r
    s.shapes.add_picture(path,Inches(bl+(bw-w)/2),Inches(bt+(bh-h)/2),Inches(w),Inches(h))

# 1. 타이틀
s=slide(NAVY2)
tb(s,1.0,2.35,11.3,0.6,[[("기말 발표",20,True,AMBER,FONT)]])
tb(s,1.0,2.95,11.5,1.4,[[("설계대로 구현하고, AI를 통제하다",48,True,WHITE,FONT)]],anchor=MSO_ANCHOR.MIDDLE)
bar(s,1.02,4.35,3.0,AMBER,0.08)
tb(s,1.0,4.6,11.3,0.9,[[("의료·자동차 보험 시스템",26,True,RGBColor(0xCF,0xE0,0xEC),FONT)]])

# 2. 세 주체
s=slide(); title(s,"세 주체")
data=[("고객",CUST,"가입·청구·납부"),
      ("직원",STAFF,"가입·지급 심사"),
      ("시스템",SYS,"계약·고지서·지급 자동화")]
xs=[1.55,5.42,9.29]
for (lab,col,role),x in zip(data,xs):
    box(s,x,1.9,2.5,2.5,col,None,[[(lab,28,True,actc(col),FONT)]],shape=MSO_SHAPE.OVAL,wrap=False)
    tb(s,x-1.0,4.75,4.5,0.8,[[(role,23,True,NAVY,FONT)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# 3. 두 가지 보험
s=slide(); title(s,"두 가지 보험")
for nm,ess,top,x in [("의료보험","아플 때 — 병원비",NAVY,1.55),("자동차보험","사고 났을 때 — 수리·치료비",AMBER,7.0)]:
    box(s,x,2.35,4.78,2.9,WHITE,LINE,[],lw=1.5)
    bar(s,x,2.35,4.78,top,0.18)
    tb(s,x,3.15,4.78,1.0,[[(nm,38,True,NAVY,FONT)]],align=PP_ALIGN.CENTER)
    tb(s,x,4.3,4.78,0.8,[[(ess,24,False,MUTED,FONT)]],align=PP_ALIGN.CENTER)

# 4. 의료보험 — 분기
s=slide(); title(s,"의료보험")
tb(s,0.7,1.65,12,0.9,[[("병원비를 보험금으로 돌려주는 보험",32,True,NAVY,FONT)]],anchor=MSO_ANCHOR.MIDDLE)
box(s,1.0,3.6,3.2,1.3,NAVY,None,[[("병원비 청구",25,True,WHITE,FONT)]])
arrow(s,4.45,2.85,30); arrow(s,4.45,4.75,30)
box(s,5.4,2.45,7.0,1.35,CUST_BG,CUST,[[("100만 원 ",22,True,NAVY,FONT),("미만",22,True,CUST,FONT),("  →  바로 입금",22,True,INK,FONT)]],align=PP_ALIGN.LEFT,lw=2.0,wrap=False)
box(s,5.4,4.35,7.0,1.35,STAFF_BG,STAFF,[[("100만 원 ",22,True,NAVY,FONT),("이상",22,True,STAFF,FONT),("  →  직원 확인 후 입금",22,True,INK,FONT)]],align=PP_ALIGN.LEFT,lw=2.0,wrap=False)

# 5. 자동차보험 — 흐름
s=slide(); title(s,"자동차보험")
tb(s,0.7,1.65,12,0.9,[[("사고 피해를 보상해 주는 보험",32,True,NAVY,FONT)]],anchor=MSO_ANCHOR.MIDDLE)
steps=["사고 접수","직원이 맡음","보상액 사정","계좌 입금"]
iw,gap=2.55,0.6; total=len(steps)*iw+(len(steps)-1)*gap; x=(SW-total)/2
for i,st in enumerate(steps):
    fill=SYS_BG if i==len(steps)-1 else CUST_BG
    box(s,x,3.55,iw,1.35,fill,CUST,[[(st,23,True,NAVY,FONT)]],lw=1.8,wrap=False)
    if i<len(steps)-1: arrow(s,x+iw+0.04,3.75)
    x+=iw+gap

# 6-8. 여정
def journey(t,steps,vtext):
    s=slide(); title(s,t); n=len(steps); iw,gap=2.45,0.7
    total=n*iw+(n-1)*gap; x=(SW-total)/2
    cmap={"cust":CUST,"staff":STAFF,"sys":SYS}; bgmap={"cust":CUST_BG,"staff":STAFF_BG,"sys":SYS_BG}
    for i,(who,lab,act) in enumerate(steps):
        box(s,x+iw/2-0.75,2.3,1.5,1.5,cmap[who],None,[[(lab,21,True,actc(cmap[who]),FONT)]],shape=MSO_SHAPE.OVAL,wrap=False)
        box(s,x,4.1,iw,1.2,bgmap[who],cmap[who],[[(act,24,True,NAVY,FONT)]],lw=2.0,wrap=False)
        if i<n-1: arrow(s,x+iw+0.08,4.35)
        x+=iw+gap
    verdict(s,vtext,t=5.75)

journey("여정 ① — 가입",
    [("cust","고객","상품 조회"),("cust","고객","가입 신청"),("staff","직원","가입 심사"),("sys","시스템","계약 생성")],
    "승인되면 계약이 바로 생깁니다.")
journey("여정 ② — 납부와 미납",
    [("cust","고객","보험료 납부"),("sys","시스템","미납 감지"),("sys","시스템","고지서 발송")],
    "납부가 밀리면 미납으로 잡히고, 매일 아침 고지서가 나갑니다 (30일 초과 시 해지 예고).")
journey("여정 ③ — 보상과 지급",
    [("cust","고객","청구·접수"),("staff","직원","보상 심사"),("sys","시스템","보험금 지급")],
    "병원비는 100만 원 미만이면 즉시, 이상이거나 자동차 사고는 직원 심사 후 지급.")

# 9. 도메인 여섯 개
s=slide(); title(s,"도메인 여섯 개")
doms=["사용자","상품","계약","청구","심사","사고이력"]; cw,gap=1.92,0.28
total=len(doms)*cw+(len(doms)-1)*gap; x=(SW-total)/2
for d in doms:
    box(s,x,3.1,cw,1.25,NAVY,None,[[(d,24,True,WHITE,FONT)]]); x+=cw+gap

# 10-16. 비교
comps=[("사용자","user","상속(User→Policyholder·Employee)과 모든 필드가 동일.","userId(String) → id(Long), birthDate Date → LocalDate."),
       ("상품","product","상품 상속과 CoverageItem 합성 관계가 동일.","productId·itemId(String) → id(Long)."),
       ("계약","contract","계약–납부–고지 합성 구조가 동일.","자동이체 정보 AutoDebit 추가, 날짜 타입 변경."),
       ("청구","claim","Claim 상속(의료/자동차)과 필드가 동일.","첨부파일 ClaimAttachment 추가, ClaimStatus 4값 → 6값(송금 결과 추적)."),
       ("__CODE__",None,None,None),
       ("심사","review","Review 상속(가입심사/지급심사) 구조가 동일.","지급심사 대상이 의료청구 → Claim(의료+자동차)으로 확대 (ADR 0009)."),
       ("사고이력","accident","사고 집계 정보(건수·지급액·면허상태)를 동일하게 보유.","외부 연동이 더미라 개별 AccidentRecord 제거, 값 객체로 단순화.")]
for name,key,same,diff in comps:
    if name=="__CODE__":
        s=slide(); title(s,"설계 = 구현 — 청구 상태값 (코드)")
        code=("public enum ClaimStatus {\n    PENDING, IN_REVIEW, APPROVED, REJECTED,   // 설계 4값\n"
              "    COMPLETED, FAILED                         // 구현 추가: 송금 성공·실패 (ADR 0007)\n}")
        box(s,1.0,2.2,11.3,2.1,RGBColor(0xF3,0xF6,0xF9),LINE,[[(l,18,False,NAVY2,MONO)] for l in code.split("\n")],align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE)
        verdict(s,"송금 결과까지 추적해야 해서 두 값을 더했습니다.",t=4.8); continue
    s=slide(); title(s,f"설계 = 구현 — {name}")
    box(s,0.6,1.5,5.9,0.55,NAVY,None,[[("설계",20,True,WHITE,FONT)]])
    box(s,6.83,1.5,5.9,0.55,AMBER,None,[[("코드 역설계",20,True,NAVY2,FONT)]])
    img_fit(s,os.path.join(CL,f"{key}-design.png"),0.6,2.15,5.9,3.2)
    img_fit(s,os.path.join(CL,f"{key}-code.png"),6.83,2.15,5.9,3.2)
    box(s,0.6,5.55,5.9,1.35,OKBG,OK,[[("= 같은 점",18,True,OK,FONT)],[(same,17,False,INK,FONT)]],align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE,lw=1.5)
    box(s,6.83,5.55,5.9,1.35,AMBERBG,AMBER,[[("≠ 다른 점",18,True,AMBER_DK,FONT)],[(diff,17,False,INK,FONT)]],align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE,lw=1.5)

# 17. 유스케이스
s=slide(); title(s,"유스케이스 — 실제 동작 경로")
img_fit(s,os.path.join(DIA,"uc-map.png"),3.0,1.55,7.3,4.3)
verdict(s,"노란 경로가 데모에서 실제로 동작합니다.",t=6.0)

# 18. AI 사용
s=slide(); title(s,"AI를 어떻게 썼는가")
box(s,0.8,1.65,5.7,1.55,OKBG,OK,[[("사람(설계자)",24,True,NAVY,FONT)],[("다이어그램 설계 · 결정(ADR) · 리뷰",20,False,INK,FONT)]],align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE,lw=1.5)
box(s,6.83,1.65,5.7,1.55,CUST_BG,CUST,[[("AI",24,True,NAVY,FONT)],[("구현 · 테스트 · 리팩토링",20,False,INK,FONT)]],align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE,lw=1.5)
steps=["설계","취조","계획","TDD 구현","리뷰"]; cw,gap=2.1,0.5; total=len(steps)*cw+(len(steps)-1)*gap; x=(SW-total)/2
for i,st in enumerate(steps):
    top=AMBER if st=="TDD 구현" else NAVY
    box(s,x,3.7,cw,1.1,RGBColor(0xEE,0xF3,0xF7),LINE,[[(st,22,True,NAVY,FONT)]],lw=1.3)
    bar(s,x,3.7,cw,top,0.12)
    if i<len(steps)-1: arrow(s,x+cw+0.02,3.9,28)
    x+=cw+gap
verdict(s,"설계와 결정은 사람이, 구현은 AI가 했습니다.",t=5.3)

# 19. AI 문제
s=slide(); title(s,"AI를 쓰며 생긴 문제와 해결")
cards=[("① 성급한 추상화","실익 없는 계층 분리 → 되돌림"),("② 명세와 불일치","자동 배정 → 수동 배정으로 정정"),
       ("③ 계층 패턴 일탈","점검으로 식별·관리"),("④ 용어·도메인 오염","용어집·ADR로 사전 차단")]
pos=[(0.8,1.65),(6.83,1.65),(0.8,3.6),(6.83,3.6)]
for (h,fix),(x,y) in zip(cards,pos):
    box(s,x,y,5.7,1.75,WHITE,LINE,[[(h,24,True,NAVY,FONT)],[(fix,20,False,OK,FONT)]],align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE,lw=1.3)
    bar2=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(0.14),Inches(1.75)); bar2.shadow.inherit=False
    bar2.fill.solid(); bar2.fill.fore_color.rgb=AMBER; bar2.line.fill.background()
verdict(s,"AI는 빠르지만, 방향은 사람이 잡았습니다.",t=5.75)

# 20. 결론
s=slide(NAVY2)
tb(s,1.0,2.6,11.3,1.0,[[("결론",40,True,WHITE,FONT)]])
tb(s,1.0,3.85,11.3,1.3,[[("설계한 그대로 구현했고, 그 과정에서 AI를 통제했습니다.",30,True,AMBER,FONT)]])

out=os.path.join(HERE,"보험발표-편집가능.pptx"); prs.save(out)
print("saved:",out,"| slides:",len(prs.slides._sldIdLst))
