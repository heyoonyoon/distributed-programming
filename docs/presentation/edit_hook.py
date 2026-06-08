#!/usr/bin/env python3
# 'SessionStart 훅' 페이지 추가 (in-place, 사용자 편집 보존: checkout 금지).
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE=os.path.dirname(os.path.abspath(__file__)); PPTX=os.path.join(HERE,"보험발표-figma.pptx")
def C(h): return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
BLACK=C("1a1a1a"); WHITE=C("ffffff"); GRAY=C("6b7280"); NAVY=C("1f1d3d"); ACCENT=C("5b5bd6"); LILAC=C("cdbff7"); CMNT=C("7c83b0")
SANS="Apple SD Gothic Neo"; MONO="Menlo"
prs=Presentation(PPTX); BLANK=prs.slide_layouts[6]
def _runs(tf,runs,align,anchor,wrap,spc):
    tf.word_wrap=wrap; tf.vertical_anchor=anchor; first=True
    for line in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False; p.alignment=align
        for (t,sz,b,c,nm) in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name=nm or SANS
            if spc is not None: r.font._rPr.set('spc',str(int(spc*100)))
def tb(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,wrap=True,spc=None):
    x=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); _runs(x.text_frame,runs,align,anchor,wrap,spc); return x
def panel(s,l,t,w,h,fill,rad=0.04):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h)); sp.shadow.inherit=False
    try: sp.adjustments[0]=rad
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.line.fill.background(); return sp
def head(s,title):
    tb(s,0.82,0.62,11.9,1.5,[[(title,44,True,BLACK,SANS)]],spc=-1.0)
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.9),Inches(1.7),Inches(0.95),Inches(0.09)); bar.shadow.inherit=False
    bar.fill.solid(); bar.fill.fore_color.rgb=ACCENT; bar.line.fill.background()
def title_of(sl):
    for sh in sl.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip(): return sh.text_frame.text.strip().split("\n")[0]
    return ""

s=prs.slides.add_slide(BLANK); f=s.background.fill; f.solid(); f.fore_color.rgb=WHITE
head(s,"훅 — 방법론을 매번 자동 주입")
panel(s,0.9,2.4,11.55,2.5,NAVY)
lines=[("// .claude 훅 설정 (SessionStart)",CMNT),
       ('"SessionStart": startup | clear | compact',WHITE),
       ("    → using-superpowers 스킬 규칙 자동 주입",LILAC),
       ("",WHITE),
       ("# 브레인스토밍 · 취조 · 계획 · TDD 규칙을",CMNT),
       ("# 세션이 열릴 때마다 자동으로 불러온다",CMNT)]
tb(s,1.3,2.65,10.9,2.05,[[(t,16,False,c,MONO)] for t,c in lines],anchor=MSO_ANCHOR.MIDDLE)
tb(s,0.9,5.2,11.6,1.0,[[("매번 '이렇게 일해'라고 시키지 않아도,",24,True,BLACK,SANS)],
    [("세션이 열릴 때마다 일하는 방식이 강제됐다.",24,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spc=-0.4)
new=s

RID=qn('r:id'); sldIdLst=prs.slides._sldIdLst
id2el={int(el.get('id')):el for el in list(sldIdLst)}
el=id2el[new.slide_id]; sldIdLst.remove(el)
pos=len(list(sldIdLst))
for i,e in enumerate(list(sldIdLst)):
    sl=next(x for x in prs.slides if x.slide_id==int(e.get('id')))
    if title_of(sl).startswith("개발 파이프라인"): pos=i+1; break
sldIdLst.insert(pos,el)
prs.save(PPTX)
print("훅 페이지 삽입 | total:", len(prs.slides._sldIdLst), "| pos:", pos)
