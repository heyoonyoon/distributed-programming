#!/usr/bin/env python3
# 개발 파이프라인 / TDD / Red-Green-Refactor 3페이지를 '사람과 AI 분업' 뒤에 삽입 (in-place).
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE=os.path.dirname(os.path.abspath(__file__)); PPTX=os.path.join(HERE,"보험발표-figma.pptx")
def C(h): return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
BLACK=C("1a1a1a"); WHITE=C("ffffff"); HAIR=C("e6e6e6"); SOFT=C("f4f4f6"); GRAY=C("6b7280")
NAVY=C("1f1d3d"); ACCENT=C("5b5bd6"); ACCENT_SOFT=C("ecebfb"); WARM=C("f7ede6")
RED=C("d8392b"); GREEN=C("2e9e5b"); BLUE=C("2f6fd8")
SANS="Apple SD Gothic Neo"; MONO="Menlo"; SW=13.333
def fg(fill): return WHITE if fill in (ACCENT,NAVY) else BLACK
prs=Presentation(PPTX); BLANK=prs.slide_layouts[6]
def newslide(bg=WHITE):
    s=prs.slides.add_slide(BLANK); f=s.background.fill; f.solid(); f.fore_color.rgb=bg; return s
def _runs(tf,runs,align,anchor,wrap,spc):
    tf.word_wrap=wrap; tf.vertical_anchor=anchor; first=True
    for line in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False; p.alignment=align
        for (t,sz,b,c,nm) in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name=nm or SANS
            if spc is not None: r.font._rPr.set('spc',str(int(spc*100)))
def tb(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,wrap=True,spc=None):
    x=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); _runs(x.text_frame,runs,align,anchor,wrap,spc); return x
def box(s,l,t,w,h,fill,runs,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,rad=0.12,line=None,wrap=True,spc=None):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h)); sp.shadow.inherit=False
    try: sp.adjustments[0]=rad
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1.3)
    tf=sp.text_frame; tf.margin_top=Pt(6); tf.margin_bottom=Pt(6); tf.margin_left=Pt(10); tf.margin_right=Pt(10)
    _runs(tf,runs,align,anchor,wrap,spc); return sp
def head(s,title):
    tb(s,0.82,0.62,11.9,1.5,[[(title,46,True,BLACK,SANS)]],spc=-1.0)
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.9),Inches(1.7),Inches(0.95),Inches(0.09)); bar.shadow.inherit=False
    bar.fill.solid(); bar.fill.fore_color.rgb=ACCENT; bar.line.fill.background()
def arrow(s,l,t,sz=26,col=GRAY):
    tb(s,l,t,0.7,0.6,[[("→",sz,True,col,SANS)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

made=[]
# 1) 개발 파이프라인
s=newslide(); head(s,"개발 파이프라인"); made.append(s)
stages=[("발산","brainstorming"),("취조","grill-with-docs"),("계획","writing-plans"),("구현","executing-plans"),("검증","Codex 리뷰")]
n=len(stages); bw,gap=2.05,0.32; total=n*bw+(n-1)*gap; x=(SW-total)/2
for i,(ko,tool) in enumerate(stages):
    col=ACCENT_SOFT if ko in ("취조","구현") else SOFT
    box(s,x,3.1,bw,1.7,col,[[(ko,22,True,BLACK,SANS)],[("",6,False,GRAY,SANS)],[(tool,12,False,GRAY,MONO)]],line=HAIR,wrap=True)
    if i<n-1: arrow(s,x+bw+(gap-0.7)/2,3.65)
    x+=bw+gap
tb(s,0.9,5.5,11.6,0.9,[[("각 단계가 산출물을 남긴다 — 설계 · 용어 · 계획 · 코드 · 리뷰.",24,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,spc=-0.4)

# 2) TDD 철칙
s=newslide(); head(s,"테스트 주도 개발 (TDD)"); made.append(s)
tb(s,0.9,2.9,11.6,1.7,[[("실패하는 테스트 없이는",36,True,BLACK,SANS)],[("프로덕션 코드를 쓰지 않는다.",36,True,ACCENT,SANS)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spc=-0.6)
tb(s,0.9,5.0,11.6,1.0,[[("AI가 쏟아내는 코드를 검증 없이 통과시키지 않는 장치.",22,False,GRAY,SANS)]],align=PP_ALIGN.CENTER,spc=-0.3)

# 3) Red → Green → Refactor
s=newslide(); head(s,"Red → Green → Refactor"); made.append(s)
cyc=[("RED",RED,"실패 테스트를\n먼저 작성"),("GREEN",GREEN,"통과시킬\n최소한만 구현"),("REFACTOR",BLUE,"녹색 유지한 채\n정리·중복 제거")]
n=len(cyc); bw,gap=3.1,0.6; total=n*bw+(n-1)*gap; x=(SW-total)/2
for i,(nm,col,desc) in enumerate(cyc):
    box(s,x,3.0,bw,1.9,col,[[(nm,22,True,WHITE,MONO)],[("",6,False,WHITE,SANS)]]+[[(ln,17,True,WHITE,SANS)] for ln in desc.split("\n")],wrap=True)
    if i<n-1: arrow(s,x+bw+(gap-0.7)/2,3.7,28,ACCENT)
    x+=bw+gap
tb(s,0.9,5.5,11.6,0.9,[[("AI가 짠 코드도 이 사이클을 강제했다.",24,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,spc=-0.4)

# ===== '사람과 AI의 분업' 뒤로 이동 =====
RID=qn('r:id'); sldIdLst=prs.slides._sldIdLst
id2el={int(el.get('id')):el for el in list(sldIdLst)}
def title_of(sl):
    for sh in sl.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip(): return sh.text_frame.text.strip().split("\n")[0]
    return ""
anchor_el=None
for sl in prs.slides:
    if title_of(sl).startswith("사람과 AI"): anchor_el=id2el[sl.slide_id]; break
new_els=[id2el[s.slide_id] for s in made]
for el in new_els: sldIdLst.remove(el)
pos=list(sldIdLst).index(anchor_el)+1
for off,el in enumerate(new_els): sldIdLst.insert(pos+off,el)
prs.save(PPTX)
print("파이프라인/TDD 3페이지 삽입 완료 | total:", len(prs.slides._sldIdLst))
