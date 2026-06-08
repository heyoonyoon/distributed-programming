#!/usr/bin/env python3
# 기존 pptx의 '프로젝트 구조' 슬라이드만 제자리에서 다시 그린다(전체 재생성 금지 — 수동 편집 보존).
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE=os.path.dirname(os.path.abspath(__file__))
PPTX=os.path.join(HERE,"보험발표-figma.pptx")

def C(h): return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
BLACK=C("1a1a1a"); WHITE=C("ffffff"); HAIR=C("e6e6e6"); SOFT=C("f4f4f6"); GRAY=C("6b7280")
NAVY=C("1f1d3d"); ACCENT=C("5b5bd6"); ACCENT_SOFT=C("ecebfb"); WARM=C("f7ede6")
SANS="Apple SD Gothic Neo"; MONO="Menlo"
def fg(fill): return WHITE if fill in (ACCENT,NAVY) else BLACK

def _runs(tf,runs,align,anchor,wrap,spc):
    tf.word_wrap=wrap; tf.vertical_anchor=anchor; first=True
    for line in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False; p.alignment=align
        for (t,sz,b,c,nm) in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name=nm or SANS
            if spc is not None: r.font._rPr.set('spc',str(int(spc*100)))

def tb(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,wrap=True,spc=None):
    x=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); _runs(x.text_frame,runs,align,anchor,wrap,spc); return x

def box(s,l,t,w,h,fill,runs,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,rad=0.16,line=None,wrap=True,spc=None):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h)); sp.shadow.inherit=False
    try: sp.adjustments[0]=rad
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1.4)
    tf=sp.text_frame; tf.margin_top=Pt(5); tf.margin_bottom=Pt(5); tf.margin_left=Pt(12); tf.margin_right=Pt(12)
    _runs(tf,runs,align,anchor,wrap,spc); return sp

def head(s,title,ink=BLACK):
    tb(s,0.82,0.62,11.9,1.5,[[(title,52,True,ink,SANS)]],spc=-1.0)
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.9),Inches(1.78),Inches(0.95),Inches(0.09)); bar.shadow.inherit=False
    bar.fill.solid(); bar.fill.fore_color.rgb=ACCENT; bar.line.fill.background()

def varrow(s,l,t):
    tb(s,l,t,0.9,0.32,[[("↓",18,True,GRAY,SANS)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

def clear(slide):
    for sp in list(slide.shapes):
        sp._element.getparent().remove(sp._element)

prs=Presentation(PPTX)
# '프로젝트 구조' 제목을 가진 슬라이드 찾기
target=None
for sl in prs.slides:
    for sh in sl.shapes:
        if sh.has_text_frame and "프로젝트 구조" in sh.text_frame.text:
            target=sl; break
    if target: break
assert target is not None, "프로젝트 구조 슬라이드를 찾지 못함"
clear(target); s=target

head(s,"프로젝트 구조")

# 백엔드 — 계층형
tb(s,0.85,2.0,3.4,0.5,[[("백엔드 · Spring",20,True,ACCENT,SANS)],[("계층형",15,False,GRAY,SANS)]],spc=-0.3)
blayers=[("controller","요청 받기·응답"),("service","업무 흐름·규칙"),
         ("domain","엔티티·도메인 규칙"),("repository","JPA 저장·조회"),("DB","")]
y=2.95
for i,(nm,desc) in enumerate(blayers):
    fillc=ACCENT if nm=="DB" else (ACCENT_SOFT if nm=="domain" else SOFT)
    runs=[[(nm,17,True,fg(fillc) if nm=="DB" else BLACK,MONO)]]
    if desc: runs[0].append(("   "+desc,12,False,GRAY,SANS))
    box(s,0.85,y,3.45,0.6,fillc,runs,align=PP_ALIGN.LEFT if desc else PP_ALIGN.CENTER,line=HAIR,wrap=False)
    if i<len(blayers)-1: varrow(s,2.35,y+0.55)
    y+=0.78

# 프론트 — 기능형 (화면→hook→api 3단)
tb(s,4.75,2.0,3.4,0.5,[[("프론트 · React",20,True,C("b5793a"),SANS)],[("기능형",15,False,GRAY,SANS)]],spc=-0.3)
flayers=[("화면","components — 보이는 UI"),("hook","상태 관리·호출 트리거"),("api","서버에 요청 보냄")]
y=2.95
for i,(nm,desc) in enumerate(flayers):
    box(s,4.75,y,3.45,0.78,WARM,[[(nm+"   ",17,True,BLACK,MONO),(desc,13,False,GRAY,SANS)]],align=PP_ALIGN.LEFT,line=HAIR,wrap=False)
    if i<len(flayers)-1: varrow(s,6.25,y+0.72)
    y+=0.96
tb(s,4.75,y+0.0,3.6,0.9,[[("한 기능 = 이 셋이 한 폴더",13,True,BLACK,SANS)],
    [("features/ : auth·contracts·claims",12,False,GRAY,MONO)],
    [("benefit-review·profile … ",12,False,GRAY,MONO)]],spc=-0.2)

# DTO — 큰 박스로 승격 (양쪽을 잇는 데이터)
box(s,8.65,2.95,3.9,2.55,NAVY,[[("DTO",30,True,WHITE,MONO)],[("",10,False,WHITE,SANS)],
    [("화면(api)과 서버(controller)가",16,True,WHITE,SANS)],[("주고받는 데이터 형식",16,True,WHITE,SANS)],
    [("",8,False,WHITE,SANS)],
    [("요청 DTO  →  서버로 보낼 값",14,False,C("cdbff7"),SANS)],
    [("응답 DTO  ←  화면에 뿌릴 값",14,False,C("cdbff7"),SANS)]],align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE,rad=0.06)
tb(s,8.65,5.6,3.9,0.5,[[("엔티티를 그대로 노출 안 하고 DTO로 변환",12,False,GRAY,SANS)]],spc=-0.2,align=PP_ALIGN.CENTER)

tb(s,0.85,6.55,11.6,0.7,[[("백엔드는 계층, 프론트는 기능 단위 — 둘 사이는 DTO로 주고받는다.",24,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,spc=-0.4)

prs.save(PPTX)
print("updated 프로젝트 구조 in place:", PPTX, "| total slides:", len(prs.slides._sldIdLst))
