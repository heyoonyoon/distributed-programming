#!/usr/bin/env python3
# cmux 터미널 이미지 추가 + 비용 페이지에 5시간 한도% + AI 문제 페이지 추가 (in-place).
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

HERE=os.path.dirname(os.path.abspath(__file__)); PPTX=os.path.join(HERE,"보험발표-figma.pptx")
DIA=os.path.join(HERE,"diagrams"); SHOT=os.path.join(DIA,"pdfimg","p5_0.png")
def C(h): return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
BLACK=C("1a1a1a"); WHITE=C("ffffff"); HAIR=C("e6e6e6"); SOFT=C("f4f4f6"); GRAY=C("6b7280")
NAVY=C("1f1d3d"); ACCENT=C("5b5bd6"); ACCENT_SOFT=C("ecebfb"); WARM=C("f7ede6")
SANS="Apple SD Gothic Neo"; MONO="Menlo"; SW=13.333
prs=Presentation(PPTX); BLANK=prs.slide_layouts[6]
def newslide(bg=WHITE):
    s=prs.slides.add_slide(BLANK); f=s.background.fill; f.solid(); f.fore_color.rgb=bg; return s
def _runs(tf,runs,align,anchor,wrap,spc):
    tf.word_wrap=wrap; tf.vertical_anchor=anchor; first=True
    for line in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False; p.alignment=align
        for (t,sz,b,c,nm) in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name=nm or SANS
            if spc is not None: r.font._rPr.set('spc',str(int(spc*100)))
def tb(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,wrap=True,spc=None):
    x=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); _runs(x.text_frame,runs,align,anchor,wrap,spc); return x
def box(s,l,t,w,h,fill,runs,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,rad=0.12,line=None,wrap=True,spc=None):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h)); sp.shadow.inherit=False
    try: sp.adjustments[0]=rad
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1.3)
    tf=sp.text_frame; tf.margin_top=Pt(6); tf.margin_bottom=Pt(6); tf.margin_left=Pt(16); tf.margin_right=Pt(14)
    _runs(tf,runs,align,anchor,wrap,spc); return sp
def head(s,title):
    tb(s,0.82,0.5,11.9,1.3,[[(title,44,True,BLACK,SANS)]],spc=-1.0)
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.9),Inches(1.55),Inches(0.95),Inches(0.09)); bar.shadow.inherit=False
    bar.fill.solid(); bar.fill.fore_color.rgb=ACCENT; bar.line.fill.background()
def arrow(s,l,t,sz=34,col=ACCENT):
    tb(s,l,t,0.8,0.7,[[("→",sz,True,col,SANS)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
def img_fit(s,path,bl,bt,bw,bh):
    iw,ih=Image.open(path).size; r=iw/ih; w,h=bw,bw/r
    if h>bh: h,w=bh,bh*r
    return s.shapes.add_picture(path,Inches(bl+(bw-w)/2),Inches(bt+(bh-h)/2),Inches(w),Inches(h))
def clear(sl):
    for sp in list(sl.shapes): sp._element.getparent().remove(sp._element)
def title_of(sl):
    for sh in sl.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip(): return sh.text_frame.text.strip().split("\n")[0]
    return ""
def find(prefix):
    for sl in prs.slides:
        if title_of(sl).startswith(prefix): return sl
    return None

# 1) cmux 페이지에 실제 터미널 이미지
s=find("cmux"); clear(s); head(s,"cmux — 두 Claude를 붙이다")
img_fit(s,SHOT,1.5,1.85,10.3,3.9)
tb(s,0.9,6.0,11.6,1.1,[[("프론트·백엔드 Claude를 한 화면에 띄우고, cmux로 서로",23,True,BLACK,SANS)],
    [("메시지를 주고받게 했다 — 한쪽이 API를 바꾸면 다른 쪽이 맞췄다.",23,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spc=-0.4)

# 2) 비용 페이지에 스케일(5시간 한도 %)
s=find("AI는 공짜"); clear(s); head(s,"AI는 공짜가 아니다")
tb(s,0.9,2.5,11.6,1.2,[[("한 Epic에  $39.63  ·  토큰 5,100만",40,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spc=-0.6)
box(s,2.4,4.3,3.9,1.4,ACCENT,[[("5시간 한도",16,True,WHITE,SANS)],[("96% 소진",30,True,WHITE,SANS)]],anchor=MSO_ANCHOR.MIDDLE)
box(s,7.0,4.3,3.9,1.4,ACCENT_SOFT,[[("주간 한도",16,True,ACCENT,SANS)],[("27% 소진",30,True,BLACK,SANS)]],anchor=MSO_ANCHOR.MIDDLE)
tb(s,0.9,6.0,11.6,0.7,[[("Claude Code 한 세션이 5시간 토큰 한도를 거의 다 먹었다.",22,False,GRAY,SANS)]],align=PP_ALIGN.CENTER,spc=-0.3)

# 3) AI 문제 추가: 권한·git 일탈
s=newslide(); head(s,"AI가 선을 넘다")
box(s,0.85,2.9,5.25,2.2,SOFT,[[("AI가 한 것",16,True,GRAY,SANS)],[("",6,False,GRAY,SANS)],
    [("git add를 넓게 잡아 frontend 전체를",19,True,BLACK,SANS)],[("엉뚱한 PR에 섞어 커밋했다",19,True,BLACK,SANS)]],align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE,line=HAIR)
arrow(s,6.27,3.7)
box(s,7.25,2.9,5.25,2.2,ACCENT_SOFT,[[("내가 통제",16,True,ACCENT,SANS)],[("",6,False,GRAY,SANS)],
    [("점검으로 발견 + '커밋은 사용자',",19,True,BLACK,SANS)],[("요청 시에만' 규칙으로 경계 관리",19,True,BLACK,SANS)]],align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE)
tb(s,0.9,5.5,11.6,0.8,[[("AI는 시키지 않은 범위까지 건드린다 — 권한 경계를 사람이 정해야 한다.",22,False,GRAY,SANS)]],align=PP_ALIGN.CENTER,spc=-0.3)
new_problem=s

# 위치: '왜 Codex' 뒤(문제 묶음)로 이동
RID=qn('r:id'); sldIdLst=prs.slides._sldIdLst
id2el={int(el.get('id')):el for el in list(sldIdLst)}
el=id2el[new_problem.slide_id]; sldIdLst.remove(el)
def idx_after(prefix):
    for i,e in enumerate(list(sldIdLst)):
        sl=next(x for x in prs.slides if x.slide_id==int(e.get('id')))
        if title_of(sl).startswith(prefix): return i+1
    return len(list(sldIdLst))
sldIdLst.insert(idx_after("왜 Codex"),el)
prs.save(PPTX)
print("완료 | total:", len(prs.slides._sldIdLst))
