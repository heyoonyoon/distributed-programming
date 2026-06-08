#!/usr/bin/env python3
# ① '프론트 클릭' 코드 슬라이드를 화면+hook 두 파일로 정정 (in-place).
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE=os.path.dirname(os.path.abspath(__file__)); PPTX=os.path.join(HERE,"보험발표-figma.pptx")
def C(h): return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
BLACK=C("1a1a1a"); WHITE=C("ffffff"); GRAY=C("6b7280"); NAVY=C("1f1d3d"); ACCENT=C("5b5bd6")
LILAC=C("cdbff7"); CMNT=C("7c83b0"); SANS="Apple SD Gothic Neo"; MONO="Menlo"

def _runs(tf,runs,align,anchor,wrap,spc):
    tf.word_wrap=wrap; tf.vertical_anchor=anchor; first=True
    for line in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False; p.alignment=align
        for (t,sz,b,c,nm) in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name=nm or SANS
            if spc is not None: r.font._rPr.set('spc',str(int(spc*100)))
def tb(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,wrap=True,spc=None):
    x=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); _runs(x.text_frame,runs,align,anchor,wrap,spc); return x
def panel(s,l,t,w,h,fill,rad=0.04):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h)); sp.shadow.inherit=False
    try: sp.adjustments[0]=rad
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.line.fill.background(); return sp
def head(s,title):
    tb(s,0.82,0.62,11.9,1.5,[[(title,52,True,BLACK,SANS)]],spc=-1.0)
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.9),Inches(1.78),Inches(0.95),Inches(0.09)); bar.shadow.inherit=False
    bar.fill.solid(); bar.fill.fore_color.rgb=ACCENT; bar.line.fill.background()
def kw(l):
    ls=l.strip()
    if ls.startswith("//"): return CMNT
    if any(k in l for k in ("async ","function","const ","return","await","import","export")): return LILAC
    return WHITE
def clear(sl):
    for sp in list(sl.shapes): sp._element.getparent().remove(sp._element)

prs=Presentation(PPTX)
target=None
for sl in prs.slides:
    for sh in sl.shapes:
        if sh.has_text_frame and "selectProduct가 실행" in sh.text_frame.text:
            target=sl; break
    if target: break
assert target is not None, "① 슬라이드를 찾지 못함"
clear(target); s=target
head(s,"프론트 — 화면과 hook")
tb(s,0.9,1.95,11.6,0.5,[[("①   ",18,True,ACCENT,SANS),
    ("ProductsView.tsx (화면)  +  useCustomerContracts.ts (hook)",16,False,GRAY,MONO)]],spc=-0.2)
panel(s,0.9,2.5,11.55,3.45,NAVY)
lines=["// ProductsView.tsx — 화면: hook에서 받아와 버튼에 연결",
 "const { selectProduct } = state",
 "<button onClick={() => selectProduct(product.id)}>…</button>","",
 "// useCustomerContracts.ts — hook: 상태 보관 + 호출",
 "async function selectProduct(id) {",
 "  setSelectedProduct(await contractsApi.getProduct(id))","}"]
tb(s,1.3,2.7,10.9,3.05,[[(l,15,False,kw(l),MONO)] for l in lines],anchor=MSO_ANCHOR.MIDDLE)
tb(s,0.9,6.15,11.6,0.8,[[("화면의 버튼이 hook의 selectProduct를 호출한다.",24,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,spc=-0.4)
prs.save(PPTX)
print("① 정정 완료")
