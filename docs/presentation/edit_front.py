from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY=RGBColor(0x0B,0x3D,0x5C); AMBER=RGBColor(0xF2,0xA9,0x00); INK=RGBColor(0x1C,0x2B,0x36)
MUTED=RGBColor(0x5B,0x71,0x85); PANEL=RGBColor(0xEE,0xF3,0xF7); LINE=RGBColor(0xD7,0xE1,0xEA)
WHITE=RGBColor(0xFF,0xFF,0xFF); LIGHT=RGBColor(0xF5,0xF9,0xFC)
CUST=RGBColor(0x2C,0x6F,0x9C); STAFF=RGBColor(0x1F,0x9E,0x8F); SYS=RGBColor(0x0B,0x3D,0x5C)
FONT="Apple SD Gothic Neo"; SW=13.333

def clear(s):
    for sh in list(s.shapes):
        sh._element.getparent().remove(sh._element)

def _set(tf,runs,align,anchor,wrap=True):
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    first=True
    for line in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=align
        for (t,sz,b,c,nm) in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b
            r.font.color.rgb=c; r.font.name=nm or FONT

def tb(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    x=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); _set(x.text_frame,runs,align,anchor); return x

def box(s,l,t,w,h,fill,line,runs,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,lw=1.2,shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp=s.shapes.add_shape(shape,Inches(l),Inches(t),Inches(w),Inches(h)); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    tf=sp.text_frame; tf.margin_top=Pt(4); tf.margin_bottom=Pt(4); tf.margin_left=Pt(8); tf.margin_right=Pt(8)
    _set(tf,runs,align,anchor); return sp

def bar(s,l,t,w,col,h=0.12):
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h)); r.shadow.inherit=False
    r.fill.solid(); r.fill.fore_color.rgb=col; r.line.fill.background()

def title(s,text):
    tb(s,0.6,0.45,12,0.9,[[(text,30,True,NAVY,FONT)]],anchor=MSO_ANCHOR.MIDDLE)
    bar(s,0.62,1.32,2.6,AMBER,0.06)

def arrow(s,l,t,sz=30):
    tb(s,l,t,0.6,0.6,[[("→",sz,True,AMBER,FONT)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

p=Presentation("보험발표-편집가능.pptx")
S=list(p.slides)

# ===== [1] 세 주체 =====
s=S[1]; clear(s); title(s,"세 주체")
data=[("고객",CUST,"고객","보험에 가입하고\n청구·납부한다"),
      ("직원",STAFF,"보험사 직원","가입과 지급을\n심사한다"),
      ("시스템",SYS,"시스템","계약·고지서·지급을\n자동 처리한다")]
xs=[1.85,5.62,9.39]
for (lab,col,nm,role),x in zip(data,xs):
    box(s,x,1.85,2.1,2.1,col,None,[[(lab,19,True,WHITE,FONT)]],shape=MSO_SHAPE.OVAL)
    tb(s,x-0.6,4.15,3.3,0.6,[[(nm,24,True,NAVY,FONT)]],align=PP_ALIGN.CENTER)
    tb(s,x-0.6,4.8,3.3,1.1,[[(l,18,False,MUTED,FONT)] for l in role.split("\n")],align=PP_ALIGN.CENTER)

# ===== [2] 두 가지 보험 =====
s=S[2]; clear(s); title(s,"두 가지 보험")
cards=[("의료보험","아플 때 — 병원비",NAVY,1.7),("자동차보험","사고 났을 때 — 수리·치료비",AMBER,7.13)]
for nm,ess,top,x in cards:
    box(s,x,2.5,4.5,2.6,LIGHT,LINE,[],lw=1.0)
    bar(s,x,2.5,4.5,top,0.14)
    tb(s,x,3.15,4.5,0.9,[[(nm,32,True,NAVY,FONT)]],align=PP_ALIGN.CENTER)
    tb(s,x,4.15,4.5,0.7,[[(ess,20,False,MUTED,FONT)]],align=PP_ALIGN.CENTER)

# ===== [3] 의료보험 — 분기 그림 =====
s=S[3]; clear(s); title(s,"의료보험")
tb(s,0.7,1.7,12,0.8,[[("병원비를 보험금으로 돌려주는 보험",28,True,NAVY,FONT)]],anchor=MSO_ANCHOR.MIDDLE)
box(s,1.2,3.7,3.0,1.1,NAVY,None,[[("병원비 청구",22,True,WHITE,FONT)]])
arrow(s,4.4,2.95,26); arrow(s,4.4,4.55,26)
box(s,5.2,2.6,7.0,1.2,LIGHT,STAFF,[[("100만 원 ",19,True,NAVY,FONT),("미만",19,True,STAFF,FONT),("  →  바로 입금",19,True,INK,FONT)]],align=PP_ALIGN.LEFT,lw=1.5)
box(s,5.2,4.2,7.0,1.2,LIGHT,AMBER,[[("100만 원 ",19,True,NAVY,FONT),("이상",19,True,RGBColor(0xC8,0x86,0x0A),FONT),("  →  직원 확인 후 입금",19,True,INK,FONT)]],align=PP_ALIGN.LEFT,lw=1.5)

# ===== [4] 자동차보험 — 흐름 그림 =====
s=S[4]; clear(s); title(s,"자동차보험")
tb(s,0.7,1.7,12,0.8,[[("사고 피해를 보상해 주는 보험",28,True,NAVY,FONT)]],anchor=MSO_ANCHOR.MIDDLE)
steps=["사고 접수","직원이 맡음","보상액 사정","계좌 입금"]
iw,gap=2.5,0.55; total=len(steps)*iw+(len(steps)-1)*gap; x=(SW-total)/2
for i,st in enumerate(steps):
    fill=SYS if i==len(steps)-1 else LIGHT; tc=WHITE if i==len(steps)-1 else NAVY
    box(s,x,3.6,iw,1.2,fill,LINE,[[(st,20,True,tc,FONT)]],lw=1.3)
    if i<len(steps)-1: arrow(s,x+iw+0.02,3.75)
    x+=iw+gap

p.save("보험발표-편집가능.pptx")
print("done: slides", len(p.slides._sldIdLst))
