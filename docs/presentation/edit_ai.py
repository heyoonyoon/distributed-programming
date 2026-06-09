#!/usr/bin/env python3
# AI 섹션 재구성(in-place): 기존 AI 2장 삭제 → 새 AI 페이지들 추가 → 결론을 맨 뒤로.
import os, copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE=os.path.dirname(os.path.abspath(__file__)); PPTX=os.path.join(HERE,"보험발표-figma.pptx")
def C(h): return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
BLACK=C("1a1a1a"); WHITE=C("ffffff"); HAIR=C("e6e6e6"); SOFT=C("f4f4f6"); GRAY=C("6b7280")
NAVY=C("1f1d3d"); ACCENT=C("5b5bd6"); ACCENT_SOFT=C("ecebfb"); WARM=C("f7ede6"); BROWN=C("b5793a")
SANS="Apple SD Gothic Neo"; MONO="Menlo"; SW=13.333
def fg(fill): return WHITE if fill in (ACCENT,NAVY) else BLACK

prs=Presentation(PPTX); BLANK=prs.slide_layouts[6]
def newslide(bg=WHITE):
    s=prs.slides.add_slide(BLANK); f=s.background.fill; f.solid(); f.fore_color.rgb=bg; return s
def _runs(tf,runs,align,anchor,wrap,spc):
    tf.word_wrap=wrap; tf.vertical_anchor=anchor; first=True
    for line in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False; p.alignment=align
        for (t,sz,b,c,nm) in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name=nm or SANS
            if spc is not None: r.font._rPr.set('spc',str(int(spc*100)))
def tb(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,wrap=True,spc=None):
    x=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); _runs(x.text_frame,runs,align,anchor,wrap,spc); return x
def box(s,l,t,w,h,fill,runs,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,rad=0.12,line=None,wrap=True,spc=None):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h)); sp.shadow.inherit=False
    try: sp.adjustments[0]=rad
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1.3)
    tf=sp.text_frame; tf.margin_top=Pt(8); tf.margin_bottom=Pt(8); tf.margin_left=Pt(16); tf.margin_right=Pt(16)
    _runs(tf,runs,align,anchor,wrap,spc); return sp
def head(s,title,ink=BLACK):
    tb(s,0.82,0.62,11.9,1.5,[[(title,46,True,ink,SANS)]],spc=-1.0)
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.9),Inches(1.7),Inches(0.95),Inches(0.09)); bar.shadow.inherit=False
    bar.fill.solid(); bar.fill.fore_color.rgb=ACCENT; bar.line.fill.background()
def arrow(s,l,t,sz=34,col=ACCENT):
    tb(s,l,t,0.8,0.7,[[("→",sz,True,col,SANS)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

def point(title,big,sub=None,big_sz=36):
    s=newslide(); head(s,title)
    tb(s,0.9,2.9,11.6,1.7,[[(big,big_sz,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,spc=-0.6)
    if sub: tb(s,0.9,4.85,11.6,1.0,[[(sub,22,False,GRAY,SANS)]],align=PP_ALIGN.CENTER,spc=-0.3)

def casepage(title,found,fixed):
    s=newslide(); head(s,title)
    box(s,0.85,3.0,5.25,2.1,SOFT,[[("AI가 한 것",16,True,GRAY,SANS)],[("",6,False,GRAY,SANS)],[(found,21,True,BLACK,SANS)]],align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE,line=HAIR)
    arrow(s,6.27,3.7)
    box(s,7.25,3.0,5.25,2.1,ACCENT_SOFT,[[("내가 통제",16,True,ACCENT,SANS)],[("",6,False,GRAY,SANS)],[(fixed,21,True,BLACK,SANS)]],align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE)

def twobox(title,lh,lb,rh,rb,lc=ACCENT_SOFT,rc=WARM,lhc=ACCENT,rhc=BROWN):
    s=newslide(); head(s,title)
    box(s,0.85,3.0,5.25,2.1,lc,[[(lh,18,True,lhc,SANS)],[("",6,False,GRAY,SANS)],[(lb,21,True,BLACK,SANS)]],align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE)
    box(s,7.25,3.0,5.25,2.1,rc,[[(rh,18,True,rhc,SANS)],[("",6,False,GRAY,SANS)],[(rb,21,True,BLACK,SANS)]],align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE)

def chips(title,big,items):
    s=newslide(); head(s,title)
    tb(s,0.9,2.5,11.6,1.0,[[(big,26,True,BLACK,SANS)]],align=PP_ALIGN.CENTER,spc=-0.4)
    n=len(items); cw,gap=3.4,0.4; total=n*cw+(n-1)*gap; x=(SW-total)/2
    for it in items:
        box(s,x,3.9,cw,1.5,SOFT,[[(it,21,True,BLACK,SANS)]],line=HAIR); x+=cw+gap

# ===== 새 AI 섹션 =====
twobox("사람과 AI의 분업","사람","설계 · 결정 · 취조","AI","구현 · 테스트 · 리팩토링")
point("AI를 어떻게 통제했나","AI가 짠 설계를 다이어그램·유스케이스에\n대고 한 질문씩 캐물었다.","결정은 전부 사람이 내렸다.",big_sz=32)
casepage("취조 ① — 빠뜨린 걸 잡다","유스케이스 17개에\n'회원가입'이 없었다","회원가입을 명시적으로\n추가하기로 결정")
casepage("취조 ② — 흔들리는 용어를 잡다","같은 사람을 Policyholder·\naccount·user로 섞어 부름","CONTEXT.md에\n표준어 하나로 못박음")
casepage("취조 ③ — 과한 구현을 막다","보험료를 '맞춤 계산'까지\n만들려 함","유스케이스 범위(예시만)로\n되돌림")
casepage("취조 ④ — 멋대로 늘린 걸 막다","상태값에 CONDITIONAL을\n임의로 추가","다이어그램 3값 유지,\n조건부는 심사가 책임")
point("AI는 공짜가 아니다","한 Epic에  $39.63  ·  토큰 5,100만",big_sz=40)
point("돈 먹은 범인","코드 생성이 아니라\n'긴 컨텍스트 재읽기'","캐시 읽기만 4,580만 토큰 — 전체 비용의 90%",big_sz=34)
chips("그래서 바꾼 것","비용을 줄이려고 습관을 바꿨다.",["Epic 끝나면\n/clear","단순 작업은\nSonnet","리뷰 횟수\n최소화"])
twobox("다중 리뷰의 두 얼굴","좋은 점","버그 6개를 미리 발견","아쉬운 점","단순 작업까지 3중 리뷰\n→ 비용 과다",lc=ACCENT_SOFT,rc=SOFT,rhc=GRAY)
point("가장 큰 실수 — 발견","자동차 사고가 직원 '심사 목록'에\n끝까지 안 떴다.",big_sz=33)
point("왜? — 설계대로였다","유스케이스(UC09): 자동차 사고는\n'접수만', 심사 큐엔 안 들어감","버그가 아니라, 내가 그렇게 설계한 것",big_sz=30)
point("어떻게 데모를 살렸나","동작하는 병원비 청구로 시연하고,\n자동차의 한계는 정직하게 설명했다.",big_sz=30)

# ===== 기존 AI 2장 삭제 + 결론 맨 뒤로 =====
from pptx.oxml.ns import qn
RID=qn('r:id')
sldIdLst=prs.slides._sldIdLst
id2el={int(el.get('id')):el for el in list(sldIdLst)}
def title_of(sl):
    for sh in sl.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text.strip().split("\n")[0]
    return ""
del_titles={"AI를 어떻게 썼나","AI가 삑사리 난 곳, 어떻게 잡았나"}
concl_el=None
for sl in list(prs.slides):
    t=title_of(sl); el=id2el[sl.slide_id]
    if t in del_titles:
        prs.part.drop_rel(el.get(RID)); sldIdLst.remove(el)   # 관계+파트까지 제거
    elif t.startswith("설계한 그대로"):
        concl_el=el
# 결론을 맨 끝으로 이동
if concl_el is not None:
    sldIdLst.remove(concl_el); sldIdLst.append(concl_el)

prs.save(PPTX)
print("AI 섹션 재구성 완료 | total slides:", len(prs.slides._sldIdLst))
